# -*- coding: utf-8 -*-
from typing import Optional, List, Dict
from pptx import Presentation
import os

from .base_reader import BaseReader
from util.logging.logger import get_logger

logger = get_logger(__name__)

class PptxReader(BaseReader):
    """PowerPoint文件读取器类，支持.pptx格式"""
    
    def read_text(self, file_path: str) -> Optional[str]:
        """
        读取PowerPoint文件中的文本内容
        
        Args:
            file_path: PowerPoint文件路径
            
        Returns:
            str: PowerPoint文件中的文本内容，如果读取失败则返回None
        """
        try:
            if not self._check_file_exists(file_path):
                return None
            
            # 打开演示文稿
            prs = Presentation(file_path)
            text_content = []
            
            # 遍历所有幻灯片
            for i, slide in enumerate(prs.slides, 1):
                # 添加幻灯片标题
                text_content.append(f"\n=== Slide {i} ===\n")
                
                # 读取幻灯片中的所有形状
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())
                
                # 读取表格内容
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        # 处理表格内容
                        for row in table.rows:
                            row_cells = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                if cell_text:  # 只添加非空单元格
                                    row_cells.append(cell_text)
                            if row_cells:  # 只添加非空行
                                text_content.append(" | ".join(row_cells))
            
            # 合并所有文本内容
            full_text = "\n".join(text_content)
            
            if full_text.strip():
                logger.info(f"成功读取PowerPoint文件文本内容: {file_path}")
                return full_text
            else:
                logger.warning(f"PowerPoint文件内容为空: {file_path}")
                return None
                
        except Exception as e:
            logger.error(f"读取PowerPoint文件文本内容失败: {str(e)}")
            return None
    
    def get_slide_count(self, file_path: str) -> Optional[int]:
        """
        获取PowerPoint文件中的幻灯片数量
        
        Args:
            file_path: PowerPoint文件路径
            
        Returns:
            int: 幻灯片数量，如果读取失败则返回None
        """
        try:
            if not self._check_file_exists(file_path):
                return None
            
            prs = Presentation(file_path)
            slide_count = len(prs.slides)
            logger.info(f"成功获取PowerPoint文件幻灯片数量: {file_path} - {slide_count}")
            return slide_count
                
        except Exception as e:
            logger.error(f"获取PowerPoint文件幻灯片数量失败: {str(e)}")
            return None
    
    def read_slide(self, file_path: str, slide_index: int) -> Optional[str]:
        """
        读取指定幻灯片的文本内容
        
        Args:
            file_path: PowerPoint文件路径
            slide_index: 幻灯片索引（从0开始）
            
        Returns:
            str: 幻灯片的文本内容，如果读取失败则返回None
        """
        try:
            if not self._check_file_exists(file_path):
                return None
            
            prs = Presentation(file_path)
            
            if slide_index < 0 or slide_index >= len(prs.slides):
                logger.error(f"幻灯片索引超出范围: {slide_index}")
                return None
            
            slide = prs.slides[slide_index]
            text_content = []
            
            # 读取幻灯片中的所有形状
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
            
            # 读取表格内容
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    # 处理表格内容
                    for row in table.rows:
                        row_cells = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:  # 只添加非空单元格
                                row_cells.append(cell_text)
                        if row_cells:  # 只添加非空行
                            text_content.append(" | ".join(row_cells))
            
            # 合并文本内容
            slide_text = "\n".join(text_content)
            
            if slide_text.strip():
                logger.info(f"成功读取PowerPoint文件幻灯片内容: {file_path} - Slide {slide_index + 1}")
                return slide_text
            else:
                logger.warning(f"PowerPoint文件幻灯片内容为空: {file_path} - Slide {slide_index + 1}")
                return None
                
        except Exception as e:
            logger.error(f"读取PowerPoint文件幻灯片内容失败: {str(e)}")
            return None

# 创建全局单例实例
pptx_reader = PptxReader() 